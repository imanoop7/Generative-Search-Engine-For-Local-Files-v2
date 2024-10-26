import os
import faiss
import numpy as np
from sentence_transformers import SentenceTransformer
import PyPDF2
import docx
from pptx import Presentation
import json
import streamlit as st
import re
import ollama
from streamlit_lottie import st_lottie
import requests
from PIL import Image
import io
import base64
import fitz  # PyMuPDF for PDF image extraction
from docx import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import tempfile

print("Starting the application...")

# Add these functions here, before they are called
def extract_images_from_pdf(file_path):
    images = []
    doc = fitz.open(file_path)
    for page in doc:
        image_list = page.get_images(full=True)
        for img_index, img in enumerate(image_list):
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            image = Image.open(io.BytesIO(image_bytes))
            images.append(image)
    return images

def extract_images_from_docx(file_path):
    images = []
    doc = Document(file_path)
    for rel in doc.part.rels.values():
        if "image" in rel.target_ref:
            image_data = rel.target_part.blob
            image = Image.open(io.BytesIO(image_data))
            images.append(image)
    return images

def extract_images_from_pptx(file_path):
    images = []
    prs = Presentation(file_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_bytes = image.blob
                image = Image.open(io.BytesIO(image_bytes))
                images.append(image)
    return images

# Global variables
model = SentenceTransformer('sentence-transformers/msmarco-bert-base-dot-v5')
dimension = 768
index = faiss.IndexFlatIP(dimension)
metadata = []

print(f"Initialized model and FAISS index with dimension {dimension}")

# Document reading functions
def read_pdf(file_path):
    print(f"Reading PDF: {file_path}")
    with open(file_path, 'rb') as file:
        reader = PyPDF2.PdfReader(file)
        return ' '.join([page.extract_text() for page in reader.pages])

def read_docx(file_path):
    print(f"Reading DOCX: {file_path}")
    doc = docx.Document(file_path)
    return ' '.join([para.text for para in doc.paragraphs])

def read_pptx(file_path):
    print(f"Reading PPTX: {file_path}")
    prs = Presentation(file_path)
    return ' '.join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, 'text')])

def chunk_text(text, chunk_size=500, overlap=50):
    print(f"Chunking text of length {len(text)} with chunk size {chunk_size} and overlap {overlap}")
    words = text.split()
    chunks = []
    for i in range(0, len(words), chunk_size - overlap):
        chunk = ' '.join(words[i:i + chunk_size])
        chunks.append(chunk)
    print(f"Created {len(chunks)} chunks")
    return chunks

def process_image(image):
    print(f"Processing image")
    with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as temp_file:
        image.save(temp_file, format='PNG')
        return temp_file.name

# Indexing function
def index_documents(directory):
    print(f"Indexing documents in directory: {directory}")
    global metadata
    documents = []
    
    file_count = 0
    for root, _, files in os.walk(directory):
        for file in files:
            file_path = os.path.join(root, file)
            print(f"Processing file: {file_path}")
            content = ""
            images = []
            
            if file.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp')):
                with Image.open(file_path) as img:
                    images.append(img)
            elif file.endswith('.pdf'):
                content = read_pdf(file_path)
                images = extract_images_from_pdf(file_path)
            elif file.endswith('.docx'):
                content = read_docx(file_path)
                images = extract_images_from_docx(file_path)
            elif file.endswith('.pptx'):
                content = read_pptx(file_path)
                images = extract_images_from_pptx(file_path)
            elif file.endswith('.txt'):
                with open(file_path, 'r', encoding='utf-8') as f:
                    content = f.read()
            
            if content:
                chunks = chunk_text(content)
                for i, chunk in enumerate(chunks):
                    documents.append(chunk)
                    metadata.append({"path": file_path, "chunk_id": i, "type": "text"})
            
            for i, img in enumerate(images):
                img_content = process_image(img)
                documents.append(img_content)
                metadata.append({"path": file_path, "chunk_id": i, "type": "image"})
            file_count += 1
    
    if file_count == 0:
        print(f"No files found in the directory: {directory}")
        return
    
    print(f"Encoding {len(documents)} document chunks and images")
    embeddings = []
    for doc in documents:
        if isinstance(doc, str):
            embeddings.append(model.encode([doc])[0])
        else:
            # For images, use LLaVA to generate a description and then encode it
            image_description = llava_generate("Describe this image in detail.", doc)
            embeddings.append(model.encode([image_description])[0])
            # Delete the temporary file after processing
            os.unlink(doc)
    
    print(f"Adding embeddings to FAISS index")
    index.add(np.array(embeddings))
    
    # Save index and metadata
    print("Saving FAISS index and metadata")
    index_path = "document_index.faiss"
    metadata_path = "metadata.json"
    faiss.write_index(index, index_path)
    with open(metadata_path, "w") as f:
        json.dump(metadata, f)
    
    print(f"Indexed {len(documents)} document chunks and images.")
    print(f"Index saved to: {os.path.abspath(index_path)}")
    print(f"Metadata saved to: {os.path.abspath(metadata_path)}")
    
    # Add this check
    if not os.path.exists(metadata_path):
        print(f"Warning: metadata file was not created at {metadata_path}")
        print(f"Current directory contents: {os.listdir('.')}")

# Function to read document chunk
def read_document_chunk(file_path, chunk_id, chunk_type):
    print(f"Reading document chunk: {file_path}, chunk_id: {chunk_id}, type: {chunk_type}")
    try:
        if chunk_type == "text":
            content = ""
            if file_path.endswith('.pdf'):
                content = read_pdf(file_path)
            elif file_path.endswith('.docx'):
                content = read_docx(file_path)
            elif file_path.endswith('.pptx'):
                content = read_pptx(file_path)
            elif file_path.endswith('.txt'):
                with open(file_path, 'r', encoding='utf-8') as f:
                    content = f.read()
            
            chunks = chunk_text(content)
            return chunks[chunk_id] if chunk_id < len(chunks) else ""
        elif chunk_type == "image":
            if file_path.endswith('.pdf'):
                images = extract_images_from_pdf(file_path)
            elif file_path.endswith('.docx'):
                images = extract_images_from_docx(file_path)
            elif file_path.endswith('.pptx'):
                images = extract_images_from_pptx(file_path)
            else:
                images = [Image.open(file_path)]
            
            if chunk_id < len(images):
                return process_image(images[chunk_id])
            else:
                return None
        else:
            print(f"Unknown chunk type: {chunk_type}")
            return ""
    except Exception as e:
        print(f"Error reading document chunk: {e}")
        return ""

# Search function
def semantic_search(query, k=10, query_type='text'):
    print(f"Performing semantic search for query type: {query_type}, k={k}")
    
    if query_type == 'text':
        query_vector = model.encode([query])[0]
    elif query_type == 'image':
        return image_search(query, k)
    
    distances, indices = index.search(np.array([query_vector]), k)
    
    results = []
    for i, idx in enumerate(indices[0]):
        meta = metadata[idx]
        content_type = meta.get("type", "text")  # Default to 'text' if 'type' is missing
        content = read_document_chunk(meta["path"], meta["chunk_id"], content_type)
        results.append({
            "id": int(idx),
            "path": meta["path"],
            "content": content,
            "type": content_type,
            "score": float(distances[0][i])
        })
    
    print(f"Found {len(results)} search results")
    return results

def image_search(image_path, k=10):
    # This function will handle image-based search using LLaVA
    base64_image = encode_image_to_base64(image_path)
    prompt = "Describe this image in detail."
    
    response = ollama.generate(model='llava', prompt=prompt, images=[base64_image])
    image_description = response['response']
    
    # Now use the image description to perform a text-based search
    return semantic_search(image_description, k, query_type='text')

# Answer generation function
def generate_answer(query, context):
    print(f"Generating answer for query: '{query}'")
    prompt = f"""Answer the user's question using the documents given in the context. In the context are documents that should contain an answer. Please always reference the document ID (in square brackets, for example [0], [1]) of the document that was used to make a claim. Use as many citations and documents as necessary to answer the question.

Context:
{context}

Question: {query}

Answer: (Remember to use document references like [0], [1], etc.)"""

    print("Sending prompt to Ollama")
    response = ollama.generate(model='tinyllama', prompt=prompt)
    print("Received response from Ollama")
    print(f"Raw response: {response['response']}")
    return response['response']

def load_lottieurl(url: str):
    r = requests.get(url)
    if r.status_code != 200:
        return None
    return r.json()

def encode_image_to_base64(image_path):
    with open(image_path, "rb") as image_file:
        return base64.b64encode(image_file.read()).decode('utf-8')

def llava_generate(prompt, image_path):
    if image_path:
        base64_image = encode_image_to_base64(image_path)
        response = ollama.generate(model='llava', prompt=prompt, images=[base64_image])
    else:
        response = ollama.generate(model='llava', prompt=prompt)
    return response['response']

def generate_answer_with_image(query, context, image_path):
    base64_image = encode_image_to_base64(image_path)
    prompt = f"""Answer the user's question using the documents and image given in the context. Please reference the document ID (in square brackets) when using information from the text documents.

Context:
{context}

Question: {query}

Answer:"""

    response = ollama.generate(model='llava', prompt=prompt, images=[base64_image])
    return response['response']

# Streamlit UI
def main():
    print("Starting Streamlit UI")
    
    # Page config
    st.set_page_config(page_title="Local GenAI Search", page_icon="🔍", layout="wide")
    
    # Custom CSS
    st.markdown("""
    <style>
    .big-font {
        font-size: 48px !important;
        font-weight: bold;
        color: #1E90FF;
        text-align: center;
        margin-bottom: 20px;
    }
    .subtitle {
        font-size: 24px;
        color: #4682B4;
        text-align: center;
        margin-bottom: 30px;
    }
    .stButton>button {
        background-color: #4CAF50;
        color: white;
        font-size: 18px;
        padding: 10px 24px;
        border-radius: 12px;
        border: none;
        transition: all 0.3s;
    }
    .stButton>button:hover {
        background-color: #45a049;
        box-shadow: 0 4px 8px rgba(0,0,0,0.1);
    }
    .stTextInput>div>div>input {
        font-size: 16px;
        border-radius: 8px;
    }
    .stExpander {
        background-color: #f0f8ff;
        border-radius: 8px;
        margin-bottom: 10px;
    }
    .stExpander>div>div>div>div>div>p {
        font-size: 16px;
        color: #333;
    }
    </style>
    """, unsafe_allow_html=True)

    # Title and animation
    st.markdown('<p class="big-font">Local GenAI Search 🔍</p>', unsafe_allow_html=True)
    st.markdown('<p class="subtitle">Explore your documents with the power of AI!</p>', unsafe_allow_html=True)

    col1, col2, col3 = st.columns([1, 2, 1])
    with col2:
        lottie_url = "https://assets5.lottiefiles.com/packages/lf20_fcfjwiyb.json"
        lottie_json = load_lottieurl(lottie_url)
        st_lottie(lottie_json, height=200, key="coding")

    # Input for documents path
    documents_path = st.text_input("📁 Enter the path to your documents folder:", "Folder Path")
    if documents_path != "Folder Path" and not os.path.exists(documents_path):
        st.error(f"The specified path does not exist: {documents_path}")
    
    # Check if documents are indexed
    if not os.path.exists("document_index.faiss") or not os.path.exists("metadata.json"):
        st.warning("⚠️ Documents are not indexed or metadata is missing. Please run the indexing process.")
        if st.button("🚀 Index Documents", key="index_button"):
            with st.spinner("Indexing documents... This may take a while."):
                print(f"Indexing documents in {documents_path}")
                if os.path.exists(documents_path):
                    index_documents(documents_path)
                    if os.path.exists("document_index.faiss") and os.path.exists("metadata.json"):
                        st.success("✅ Indexing complete!")
                        st.rerun()  # Changed from st.experimental_rerun() to st.rerun()
                    else:
                        st.error("Indexing failed. Please check the application logs.")
                else:
                    st.error(f"The specified documents path does not exist: {documents_path}")

    # Load index and metadata if not already loaded
    global index, metadata
    if len(metadata) == 0:
        print("Loading FAISS index and metadata")
        index_path = "document_index.faiss"
        metadata_path = "metadata.json"
        try:
            if not os.path.exists(index_path):
                raise FileNotFoundError(f"Index file not found: {index_path}")
            if not os.path.exists(metadata_path):
                raise FileNotFoundError(f"Metadata file not found: {metadata_path}")
            
            index = faiss.read_index(index_path)
            with open(metadata_path, "r") as f:
                metadata = json.load(f)
            print(f"Loaded index with {index.ntotal} vectors and {len(metadata)} metadata entries")
        except FileNotFoundError as e:
            st.error(f"Error: {str(e)}. Please run the indexing process first.")
            st.error(f"Current working directory: {os.getcwd()}")
            st.error(f"Files in current directory: {os.listdir('.')}")
            return
        except json.JSONDecodeError:
            st.error(f"Error reading metadata file: {metadata_path}. Please run the indexing process again.")
            return
        except Exception as e:
            st.error(f"Unexpected error: {str(e)}. Please check the application logs.")
            return
    
    st.markdown("---")
    st.markdown("## 🤔 Ask a Question")
    question = st.text_input("What would you like to know about your documents?", "")

    col1, col2, col3 = st.columns([1, 1, 1])
    with col2:
        search_button = st.button("🔍 Search and Answer", key="search_button")

    if search_button:
        if question:
            with st.spinner("🕵️‍♀️ Searching and generating answer..."):
                search_results = semantic_search(question)
                context = "\n\n".join([f"{i}: {result['content']}" for i, result in enumerate(search_results)])
                answer = generate_answer(question, context)
                
                # Display answer and referenced documents
                st.markdown("### 🤖 AI Answer:")
                st.info(answer)
                
                # Display referenced documents
                st.markdown("### 📚 Referenced Documents:")
                rege = re.compile(r"\[Document\s+[0-9]+\]|\[[0-9]+\]")
                referenced_ids_raw = re.findall(r'\b\d+\b', ' '.join(rege.findall(answer)))
                referenced_ids = [int(s) for s in referenced_ids_raw]

                print(f"Raw answer: {answer}")
                print(f"Regex matches: {rege.findall(answer)}")
                print(f"Referenced IDs (raw): {referenced_ids_raw}")
                print(f"Referenced IDs: {referenced_ids}")

                if not referenced_ids:
                    st.warning("No specific document references found in the answer.")

                print(f"Displaying {len(referenced_ids)} referenced documents")
                for doc_id in referenced_ids:
                    if doc_id < len(search_results):
                        doc = search_results[doc_id]
                        with st.expander(f"📄 Document {doc_id} - {os.path.basename(doc['path'])}"):
                            st.write(doc['content'])
                            col1, col2 = st.columns([3, 1])
                            with col2:
                                with open(doc['path'], 'rb') as f:
                                    st.download_button("⬇️ Download file", f, file_name=os.path.basename(doc['path']))
                    else:
                        st.warning(f"Referenced document ID {doc_id} is out of range.")
        else:
            st.warning("⚠️ Please enter a question before clicking 'Search and Answer'.")

    st.markdown("---")
    st.markdown("### 🌟 Made with love by Anoop Maurya")

if __name__ == "__main__":
    main()
    print("Application finished")

